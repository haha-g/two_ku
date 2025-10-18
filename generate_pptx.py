# 运行前：pip install python-pptx requests Pillow
# 用法（Windows PowerShell）:
# 1) 将此脚本保存为 generate_pptx.py
# 2) 打开 PowerShell（以用户权限），运行：
#    python .\generate_pptx.py "C:\GitHub01\two_ku"
# 脚本会在指定目录下生成:
# 酶法塑料降解_PETase_环境修复与资源化_12页.pptx

import sys
import os
import io
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from PIL import Image

# ---------- 配置 ----------
OUT_FILENAME = "酶法塑料降解_PETase_环境修复与资源化_12页.pptx"
# 若需要替换图片，可把 URL 换成你本地路径或其他图片 URL（注意许可）
IMAGE_URLS = {
    "pet_structure": "https://upload.wikimedia.org/wikipedia/commons/6/6b/Polyethylene_terephthalate_structure.svg",
    "enzyme_schematic": "https://upload.wikimedia.org/wikipedia/commons/1/1b/Protein_domain_structure.svg",
    "process_flow": "https://upload.wikimedia.org/wikipedia/commons/4/4f/Recycle_process_icon.svg"
}
# ---------- 幻灯片内容（标题、要点、备注） ----------
slides = [
    ("酶法塑料降解（以PETase为例）在环境修复与资源化中的应用",
     ["汇报人：你的姓名（可后补）", "基于近5年文献与产业案例（12页）"],
     "开场：介绍主题、目的与总体结构。"),
    ("背景：塑料污染现状",
     ["全球塑料产量与废弃量持续增长", "常见难降解塑料：PET、PE、PP、PA", "传统回收与焚烧的局限"],
     "说明为何需要新技术（环境与资源循环双重需求）。"),
    ("研究对象与检索范围",
     ["技术聚焦：酶法降解（PETase）", "检索时间：近5年（2020–2025）", "关键词：PETase / 塑料降解酶 / 酶工程 / 工业化"],
     "说明检索范围与目标：总结原理、进展、案例与挑战。"),
    ("技术原理：PET 与 PETase",
     ["PET简介：广泛用于饮料瓶、纤维和包装", "PETase：来源Ideonella sakaiensis，水解PET→MHET→TPA+EG", "与MHETase协同可实现更完全分解"],
     "配示意图：PET分子与降解路径（图示）。"),
    ("蛋白质工程与酶改造进展",
     ["通过定向进化与理性设计提高活性与热稳定性", "改造以改善底物结合与催化效率", "LCC、PHL7等新型高效酶被发现"],
     "举例说明结构改造如何提升工业适应性。"),
    ("多酶体系与表达/工艺化策略",
     ["多酶协同（PETase+MHETase等）提升降解效率", "表面展示、分泌表达有利于量产", "固定化酶与连续反应器的工程化尝试"],
     "讨论酶在工程化中的实现方式与优势。"),
    ("应用流程与工程化步骤",
     ["原料预处理（粉碎/热处理）提高酶可达性", "酶催化（自由酶/固定化/菌体）", "单体回收并再聚合为 rPET"],
     "展示典型工业工作流程并指出关键参数。"),
    ("工业化案例：Carbios（与品牌合作）",
     ["Carbios推进酶法PET回收商业化", "示范工厂与品牌合作实现rPET回用", "产业意义：闭环回收、减少化石原料依赖"],
     "简述企业落地要点与市场合作模式。"),
    ("中国与其他企业/团队实践",
     ["源天生物等推进吨级产能与纺织链对接", "科研团队在酶筛选与反应器优化上取得进展"],
     "强调本地化产业化与技术差异化的挑战。"),
    ("监测与环境安全评估",
     ["监测指标：TPA/EG产率、残余塑料形貌、微塑料数量、TOC/COD", "风险：中间产物毒性、酶/微生物释放的生态影响"],
     "实地修复需同步环境风险监测与评估。"),
    ("优势、局限与挑战",
     ["优势：温和条件、可实现单体回收、环保友好", "局限：对高结晶PET或PE/PP效果差、成本与放大问题", "需结合预处理与工程设计降低成本"],
     "总结目前技术推广的主要障碍。"),
    ("未来方向与结论",
     ["AI辅助酶工程与高通量筛选", "热化学预处理 + 酶法耦合工艺", "加强LCA与监管标准，推动产业化"],
     "结论：酶法PET降解已从实验室向示范商业化过渡，但仍需跨学科协作。")
]

# ---------- 帮助函数：下载图片 ----------
def download_image(url):
    try:
        r = requests.get(url, timeout=15)
        r.raise_for_status()
        img = Image.open(io.BytesIO(r.content)).convert("RGB")
        return img
    except Exception as e:
        print(f"Warning: 无法下载图片 {url} ：{e}")
        return None

def save_temp_image(img, path):
    img.save(path, format="JPEG", quality=85)

# ---------- 生成 PPT ----------
def create_pptx(out_dir):
    prs = Presentation()
    # 选择空白模板页索引参考：0..n 但不同模板略有差异。为兼容性，用默认布局1（Title Slide）和 1（Title and Content）
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # 预下载图片
    imgs = {}
    for key, url in IMAGE_URLS.items():
        img = download_image(url)
        if img:
            tmp_path = os.path.join(out_dir, f"tmp_{key}.jpg")
            save_temp_image(img, tmp_path)
            imgs[key] = tmp_path
        else:
            imgs[key] = None

    for idx, (title, bullets, notes) in enumerate(slides):
        if idx == 0:
            slide = prs.slides.add_slide(title_slide_layout)
            title_tf = slide.shapes.title
            subtitle = slide.placeholders[1]
            title_tf.text = title
            subtitle.text = "\n".join(bullets)
            # Speaker notes
            slide.notes_slide.notes_text_frame.text = notes
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(bullets):
                p = body.add_paragraph() if i>0 else body.paragraphs[0]
                p.text = b
                p.level = 0
                p.font.size = Pt(18)
            # 插入示意图到某几页
            if idx in (3,5,6):  # 4th,6th,7th slides (0-based index)
                img_key = "pet_structure" if idx==3 else ("enzyme_schematic" if idx==5 else "process_flow")
                img_path = imgs.get(img_key)
                if img_path:
                    # 添加图片到右侧
                    left = Inches(6)
                    top = Inches(1.2)
                    height = Inches(3.5)
                    try:
                        slide.shapes.add_picture(img_path, left, top, height=height)
                    except Exception as e:
                        print(f"插入图片失败：{e}")
            # 备注
            slide.notes_slide.notes_text_frame.text = notes

    # 清理临时图片
    out_path = os.path.join(out_dir, OUT_FILENAME)
    try:
        prs.save(out_path)
        print(f"PPT 文件已生成：{out_path}")
    finally:
        for p in imgs.values():
            if p and os.path.exists(p):
                try:
                    os.remove(p)
                except:
                    pass

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python generate_pptx.py <输出目录>（例如 C:\\GitHub01\\two_ku）")
        sys.exit(1)
    out_dir = sys.argv[1]
    if not os.path.isdir(out_dir):
        try:
            os.makedirs(out_dir, exist_ok=True)
        except Exception as e:
            print(f"无法创建目录 {out_dir} ：{e}")
            sys.exit(1)
    create_pptx(out_dir)